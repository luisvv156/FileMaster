"""Extracción de texto para los formatos soportados por FileMaster.

Jerarquía de extracción por formato:
  PDF   → pdfplumber → PyMuPDF (fitz) → pdftotext CLI → OCR
  DOCX  → python-docx → XML manual (fallback)
  PPTX  → python-pptx → XML manual (fallback)
  Texto → lectura directa con detección de encoding
  Imagen → OCR (Tesseract vía OCRHandler)
"""

from __future__ import annotations

import logging
import re
import shutil
import subprocess
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from config.settings import IMAGE_EXTENSIONS
from core.models import ExtractionResult
from core.ocr_handler import OCRHandler

logger = logging.getLogger(__name__)

MAX_CHARS = 80_000  # Increased for better document analysis
MIN_CONTENT_LENGTH = 100  # Lower threshold for longer extraction
COVER_MAX_CHARS = 3_500
COVER_MAX_PAGES = 2

_TEXT_ENCODINGS = ("utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1")

# Mínimo de caracteres para considerar una extracción válida
_MIN_TEXT_LENGTH = MIN_CONTENT_LENGTH


class TextExtractor:
    """Extrae texto de documentos en múltiples formatos."""

    def __init__(self) -> None:
        self.ocr_handler = OCRHandler()
        self._pdftotext    = shutil.which("pdftotext")
        self._has_pdfplumber = self._check_import("pdfplumber")
        self._has_fitz     = self._check_import("fitz")
        self._has_docx     = self._check_import("docx")
        self._has_pptx     = self._check_import("pptx")
        self._has_openpyxl = self._check_import("openpyxl")
        self._has_xlrd     = self._check_import("xlrd")

        # Log de capacidades disponibles al iniciar
        caps = []
        if self._has_pdfplumber: caps.append("pdfplumber")
        if self._has_fitz:       caps.append("pymupdf")
        if self._pdftotext:      caps.append("pdftotext")
        if self._has_docx:       caps.append("python-docx")
        if self._has_pptx:       caps.append("python-pptx")
        if self._has_openpyxl:   caps.append("openpyxl")
        if self._has_xlrd:       caps.append("xlrd")
        logger.info("TextExtractor iniciado — capacidades: %s", ", ".join(caps) or "básico")

    # ------------------------------------------------------------------
    # Punto de entrada principal
    # ------------------------------------------------------------------

    def extract(self, file_path: Path) -> ExtractionResult:
        if not file_path.exists():
            return ExtractionResult("", "error", f"Archivo no encontrado: {file_path}")

        suffix = file_path.suffix.lower()

        extractors = {
            ".pdf":  self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".xlsx": self._extract_excel,
            ".xlsm": self._extract_excel,
            ".xltx": self._extract_excel,
            ".xltm": self._extract_excel,
            ".xls":  self._extract_excel,
            ".doc":  self._extract_doc_fallback,
        }

        plain_text_extensions = {
            ".txt", ".md", ".csv", ".json", ".log",
            ".py", ".js", ".html", ".xml", ".rst",
        }

        if suffix in plain_text_extensions:
            result = self._extract_plain_text(file_path)
        elif suffix in extractors:
            result = extractors[suffix](file_path)
        elif suffix in IMAGE_EXTENSIONS:
            result = self.ocr_handler.extract_text(file_path)
        else:
            return ExtractionResult("", "unsupported", f"Formato no soportado: {suffix}")

        if len(result.text) > MAX_CHARS:
            result = ExtractionResult(result.text[:MAX_CHARS], result.method, result.note)

        return result

    def extract_cover(
        self,
        file_path: Path,
        *,
        max_chars: int = COVER_MAX_CHARS,
        max_pages: int = COVER_MAX_PAGES,
    ) -> ExtractionResult:
        """Extrae texto de portada/primeras páginas para clasificación rápida.

        Este método minimiza CPU/RAM al leer solo el inicio del archivo y
        prioriza los datos de materia/docente que suelen venir en portada.
        """
        if not file_path.exists():
            return ExtractionResult("", "error", f"Archivo no encontrado: {file_path}")

        suffix = file_path.suffix.lower()

        if suffix == ".pdf":
            result = self._extract_pdf_cover(file_path, max_pages=max_pages)
        elif suffix == ".docx":
            result = self._extract_docx_cover(file_path)
        elif suffix == ".pptx":
            result = self._extract_pptx_cover(file_path, max_pages=max_pages)
        elif suffix in {".xlsx", ".xlsm", ".xltx", ".xltm", ".xls"}:
            result = self._extract_excel_cover(file_path, max_rows=120)
        elif suffix in IMAGE_EXTENSIONS:
            result = self.ocr_handler.extract_text(file_path)
        else:
            result = self.extract(file_path)

        text = (result.text or "").strip()
        if len(text) > max_chars:
            text = text[:max_chars]
        return ExtractionResult(text, f"{result.method}_cover", result.note)

    # ------------------------------------------------------------------
    # PDF — jerarquía de extractores
    # ------------------------------------------------------------------

    def _extract_pdf(self, file_path: Path) -> ExtractionResult:
        """Intenta extraer texto de un PDF con múltiples métodos en orden de calidad."""

        # Intento 1: pdfplumber — mejor para PDFs académicos con columnas y tablas
        if self._has_pdfplumber:
            result = self._extract_pdf_pdfplumber(file_path)
            if len(result.text.strip()) >= _MIN_TEXT_LENGTH:
                return result
            logger.debug("pdfplumber: texto insuficiente en %s (%d chars)", file_path.name, len(result.text))

        # Intento 2: PyMuPDF — rápido y preciso para PDFs normales
        if self._has_fitz:
            result = self._extract_pdf_fitz(file_path)
            if len(result.text.strip()) >= _MIN_TEXT_LENGTH:
                return result
            logger.debug("PyMuPDF: texto insuficiente en %s (%d chars)", file_path.name, len(result.text))

        # Intento 3: pdftotext CLI
        if self._pdftotext:
            result = self._extract_pdf_pdftotext(file_path)
            if len(result.text.strip()) >= _MIN_TEXT_LENGTH:
                return result
            logger.debug("pdftotext: texto insuficiente en %s (%d chars)", file_path.name, len(result.text))

        # Intento 4: OCR — PDF escaneado o con texto como imagen
        if self.ocr_handler.available:
            logger.info("PDF sin texto extraíble, intentando OCR: %s", file_path.name)
            result = self.ocr_handler.extract_text(file_path)
            if len(result.text.strip()) >= _MIN_TEXT_LENGTH:
                return result

        # Último recurso: regex básico (puede producir basura pero es mejor que nada)
        logger.warning(
            "Todos los métodos fallaron para '%s' — usando regex básico. "
            "El clustering puede ser impreciso para este archivo.",
            file_path.name,
        )
        return self._extract_pdf_regex(file_path)

    def _extract_pdf_cover(self, file_path: Path, *, max_pages: int) -> ExtractionResult:
        """Versión rápida para portada de PDF."""
        if self._has_fitz:
            try:
                import fitz
                pages: list[str] = []
                with fitz.open(str(file_path)) as doc:
                    for index in range(min(max_pages, len(doc))):
                        text = doc[index].get_text("text")
                        if text and text.strip():
                            pages.append(text)
                if pages:
                    return ExtractionResult("\n\n".join(pages), "pdf_fitz", "")
            except Exception as exc:
                logger.debug("PyMuPDF portada falló en %s: %s", file_path.name, exc)

        if self._has_pdfplumber:
            try:
                import pdfplumber
                pages: list[str] = []
                with pdfplumber.open(str(file_path)) as pdf:
                    for page in pdf.pages[:max_pages]:
                        text = page.extract_text()
                        if text and text.strip():
                            pages.append(text)
                if pages:
                    return ExtractionResult("\n\n".join(pages), "pdf_pdfplumber", "")
            except Exception as exc:
                logger.debug("pdfplumber portada falló en %s: %s", file_path.name, exc)

        if self._pdftotext:
            try:
                result = subprocess.run(
                    [self._pdftotext, "-f", "1", "-l", str(max_pages), "-layout", str(file_path), "-"],
                    check=False,
                    capture_output=True,
                    text=True,
                    timeout=20,
                )
                if result.stdout.strip():
                    return ExtractionResult(result.stdout.strip(), "pdf_pdftotext", "")
            except (OSError, subprocess.SubprocessError) as exc:
                logger.debug("pdftotext portada falló en %s: %s", file_path.name, exc)

        return self._extract_pdf(file_path)

    def _extract_pdf_pdfplumber(self, file_path: Path) -> ExtractionResult:
        """Extrae texto con pdfplumber — excelente para PDFs académicos."""
        try:
            import pdfplumber
            pages: list[str] = []
            with pdfplumber.open(str(file_path)) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append(text)
            return ExtractionResult("\n\n".join(pages), "pdf_pdfplumber", "")
        except Exception as exc:
            logger.warning("pdfplumber falló en %s: %s", file_path.name, exc)
            return ExtractionResult("", "pdf_pdfplumber", str(exc))

    def _extract_pdf_fitz(self, file_path: Path) -> ExtractionResult:
        """Extrae texto con PyMuPDF."""
        try:
            import fitz
            pages: list[str] = []
            with fitz.open(str(file_path)) as doc:
                for page in doc:
                    text = page.get_text("text")
                    if text.strip():
                        pages.append(text)
            return ExtractionResult("\n\n".join(pages), "pdf_fitz", "")
        except Exception as exc:
            logger.warning("PyMuPDF falló en %s: %s", file_path.name, exc)
            return ExtractionResult("", "pdf_fitz", str(exc))

    def _extract_pdf_pdftotext(self, file_path: Path) -> ExtractionResult:
        """Extrae texto con pdftotext CLI."""
        try:
            result = subprocess.run(
                [self._pdftotext, "-layout", str(file_path), "-"],
                check=False,
                capture_output=True,
                text=True,
                timeout=30,
            )
            return ExtractionResult(result.stdout.strip(), "pdf_pdftotext", "")
        except (OSError, subprocess.SubprocessError) as exc:
            logger.warning("pdftotext falló en %s: %s", file_path.name, exc)
            return ExtractionResult("", "pdf_pdftotext", str(exc))

    def _extract_pdf_regex(self, file_path: Path) -> ExtractionResult:
        """Extracción básica por regex — último recurso, puede producir ruido."""
        try:
            raw = file_path.read_bytes().decode("latin-1", errors="ignore")
            matches = re.findall(r"\(([^()]{4,200})\)", raw)
            cleaned = " ".join(
                m for m in matches
                if "/" not in m
                and "\\" not in m
                and len(m.split()) > 1
                # FIX: filtrar cadenas sin vocales (basura codificada)
                and any(c in "aeiouáéíóúAEIOUÁÉÍÓÚ" for c in m)
            )
            note = "Extracción básica — calidad limitada" if cleaned else "Sin texto extraíble"
            return ExtractionResult(cleaned, "pdf_regex", note)
        except OSError as exc:
            return ExtractionResult("", "pdf_regex", str(exc))

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _extract_docx(self, file_path: Path) -> ExtractionResult:
        if self._has_docx:
            try:
                from docx import Document
                doc = Document(str(file_path))
                parts: list[str] = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                parts.append(cell.text)
                if parts:
                    return ExtractionResult("\n".join(parts), "docx_python", "")
            except Exception as exc:
                logger.warning("python-docx falló en %s: %s", file_path.name, exc)

        return self._extract_docx_xml(file_path)

    def _extract_docx_cover(self, file_path: Path) -> ExtractionResult:
        """Extrae solo los primeros bloques de un DOCX."""
        if self._has_docx:
            try:
                from docx import Document
                doc = Document(str(file_path))
                parts: list[str] = []
                for para in doc.paragraphs[:30]:
                    if para.text.strip():
                        parts.append(para.text)
                for table in doc.tables[:2]:
                    for row in table.rows[:5]:
                        for cell in row.cells[:4]:
                            if cell.text.strip():
                                parts.append(cell.text)
                if parts:
                    return ExtractionResult("\n".join(parts), "docx_python", "")
            except Exception as exc:
                logger.debug("DOCX portada falló en %s: %s", file_path.name, exc)
        return self._extract_docx(file_path)

    def _extract_docx_xml(self, file_path: Path) -> ExtractionResult:
        try:
            with ZipFile(file_path) as archive:
                xml_data = archive.read("word/document.xml")
        except OSError:
            return ExtractionResult("", "docx", "No fue posible abrir el DOCX")
        except KeyError:
            return ExtractionResult("", "docx", "El DOCX no contiene document.xml")

        root = ET.fromstring(xml_data)
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        paragraphs: list[str] = []
        for para in root.findall(".//w:p", ns):
            texts = [n.text for n in para.findall(".//w:t", ns) if n.text]
            if texts:
                paragraphs.append("".join(texts))
        return ExtractionResult("\n".join(paragraphs), "docx_xml", "")

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def _extract_pptx(self, file_path: Path) -> ExtractionResult:
        if self._has_pptx:
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                slides: list[str] = []
                for slide in prs.slides:
                    slide_texts: list[str] = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    slide_texts.append(text)
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            slide_texts.append(f"[Notas: {notes}]")
                    if slide_texts:
                        slides.append(" ".join(slide_texts))
                if slides:
                    return ExtractionResult("\n".join(slides), "pptx_python", "")
            except Exception as exc:
                logger.warning("python-pptx falló en %s: %s", file_path.name, exc)

        return self._extract_pptx_xml(file_path)

    def _extract_pptx_cover(self, file_path: Path, *, max_pages: int) -> ExtractionResult:
        """Extrae solo las primeras diapositivas."""
        if self._has_pptx:
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                slides: list[str] = []
                limit = min(max_pages, len(prs.slides))
                for idx in range(limit):
                    slide = prs.slides[idx]
                    slide_texts: list[str] = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    slide_texts.append(text)
                    if slide_texts:
                        slides.append(" ".join(slide_texts))
                if slides:
                    return ExtractionResult("\n".join(slides), "pptx_python", "")
            except Exception as exc:
                logger.debug("PPTX portada falló en %s: %s", file_path.name, exc)
        return self._extract_pptx(file_path)

    def _extract_pptx_xml(self, file_path: Path) -> ExtractionResult:
        try:
            with ZipFile(file_path) as archive:
                slide_names = sorted(
                    n for n in archive.namelist()
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml")
                )
                slides: list[str] = []
                for name in slide_names:
                    root = ET.fromstring(archive.read(name))
                    texts = [n.text for n in root.findall(".//{*}t") if n.text]
                    if texts:
                        slides.append(" ".join(texts))
        except OSError as exc:
            return ExtractionResult("", "pptx", str(exc))
        return ExtractionResult("\n".join(slides), "pptx_xml", "")

    # ------------------------------------------------------------------
    # Excel
    # ------------------------------------------------------------------

    def _extract_excel(self, file_path: Path) -> ExtractionResult:
        suffix = file_path.suffix.lower()
        if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
            return self._extract_excel_openpyxl(file_path, max_rows=None)
        if suffix == ".xls":
            return self._extract_excel_xlrd(file_path)
        return ExtractionResult("", "excel", f"Formato Excel no soportado: {suffix}")

    def _extract_excel_cover(self, file_path: Path, *, max_rows: int = 120) -> ExtractionResult:
        suffix = file_path.suffix.lower()
        if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
            return self._extract_excel_openpyxl(file_path, max_rows=max_rows)
        if suffix == ".xls":
            return self._extract_excel_xlrd(file_path, max_rows=max_rows)
        return self._extract_excel(file_path)

    def _extract_excel_openpyxl(self, file_path: Path, max_rows: int | None) -> ExtractionResult:
        if not self._has_openpyxl:
            return ExtractionResult("", "excel_openpyxl", "openpyxl no está instalado")
        try:
            from openpyxl import load_workbook

            wb = load_workbook(filename=str(file_path), read_only=True, data_only=True)
            lines: list[str] = []
            rows_read = 0
            max_cols = 20
            for sheet in wb.worksheets:
                lines.append(f"[Hoja: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    values = []
                    for value in row[:max_cols]:
                        if value is None:
                            continue
                        text = str(value).strip()
                        if text:
                            values.append(text)
                    if values:
                        lines.append(" | ".join(values))
                    rows_read += 1
                    if max_rows is not None and rows_read >= max_rows:
                        break
                if max_rows is not None and rows_read >= max_rows:
                    break
            wb.close()
            if lines:
                return ExtractionResult("\n".join(lines), "excel_openpyxl", "")
            return ExtractionResult("", "excel_openpyxl", "No se encontró contenido legible en Excel")
        except Exception as exc:
            logger.warning("openpyxl falló en %s: %s", file_path.name, exc)
            return ExtractionResult("", "excel_openpyxl", str(exc))

    def _extract_excel_xlrd(self, file_path: Path, max_rows: int | None = None) -> ExtractionResult:
        if not self._has_xlrd:
            return ExtractionResult(
                "",
                "excel_xlrd",
                "Formato .xls requiere xlrd instalado o convertir a .xlsx",
            )
        try:
            import xlrd

            book = xlrd.open_workbook(str(file_path), on_demand=True)
            lines: list[str] = []
            rows_read = 0
            max_cols = 20
            for sheet in book.sheets():
                lines.append(f"[Hoja: {sheet.name}]")
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    values = []
                    for value in row[:max_cols]:
                        text = str(value).strip()
                        if text:
                            values.append(text)
                    if values:
                        lines.append(" | ".join(values))
                    rows_read += 1
                    if max_rows is not None and rows_read >= max_rows:
                        break
                if max_rows is not None and rows_read >= max_rows:
                    break
            book.release_resources()
            if lines:
                return ExtractionResult("\n".join(lines), "excel_xlrd", "")
            return ExtractionResult("", "excel_xlrd", "No se encontró contenido legible en .xls")
        except Exception as exc:
            logger.warning("xlrd falló en %s: %s", file_path.name, exc)
            return ExtractionResult("", "excel_xlrd", str(exc))

    # ------------------------------------------------------------------
    # DOC (Word 97-2003)
    # ------------------------------------------------------------------

    def _extract_doc_fallback(self, file_path: Path) -> ExtractionResult:
        antiword = shutil.which("antiword")
        if antiword:
            try:
                result = subprocess.run(
                    [antiword, str(file_path)],
                    capture_output=True, text=True, timeout=15,
                )
                if result.stdout.strip():
                    return ExtractionResult(result.stdout, "doc_antiword", "")
            except (OSError, subprocess.SubprocessError):
                pass
        return ExtractionResult("", "doc", "Formato .doc requiere antiword instalado")

    # ------------------------------------------------------------------
    # Texto plano
    # ------------------------------------------------------------------

    def _extract_plain_text(self, file_path: Path) -> ExtractionResult:
        for encoding in _TEXT_ENCODINGS:
            try:
                text = file_path.read_text(encoding=encoding)
                return ExtractionResult(text, "plain", "")
            except UnicodeDecodeError:
                continue
            except OSError as exc:
                return ExtractionResult("", "plain", str(exc))
        return ExtractionResult("", "plain", "No fue posible decodificar el archivo")

    # ------------------------------------------------------------------
    # Utilidades
    # ------------------------------------------------------------------

    @staticmethod
    def _check_import(module: str) -> bool:
        import importlib.util
        return importlib.util.find_spec(module) is not None
